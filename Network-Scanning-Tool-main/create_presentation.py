
"""
Generate a PowerPoint presentation for the Network Scanning Project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 120)  # Dark blue
    
    # Title
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    left = Inches(0.5)
    top = Inches(3.8)
    width = Inches(9)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    # Content
    left = Inches(0.7)
    top = Inches(1.2)
    width = Inches(8.6)
    height = Inches(5.3)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_code_slide(prs, title, code_text):
    """Add a slide with code/command examples."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
    
    # Code
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.3)
    code_box = slide.shapes.add_textbox(left, top, width, height)
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    
    p = code_frame.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Background for code
    shape = slide.shapes.add_shape(1, Inches(0.4), Inches(1.1), Inches(9.2), Inches(5.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 30, 30)
    shape.line.color.rgb = RGBColor(100, 100, 100)
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    
    return slide

def create_presentation():
    """Create the Network Scanning Project presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "Network Scanning Tool", 
                    "A Python-based Network Discovery & Analysis Solution")
    
    # Slide 2: Overview
    add_content_slide(prs, "Project Overview", [
        "• Network scanning is the process of discovering and analyzing active hosts and services on a network.",
        "• Uses multiple techniques: ICMP (ping), TCP port scanning, and ARP (Address Resolution Protocol).",
        "• Essential for network administration, security auditing, and penetration testing."
    ])
    
    # Slide 3: Features Part 1
    add_content_slide(prs, "Key Features (Part 1)", [
        "ICMP Scan: Checks if a host is up using ping (echo request)",
        "• Reports host status: UP or DOWN",
        "• Includes target IP in results",
        "",
        "TCP Port Scan: Identifies open, closed, or filtered ports",
        "• Supports custom port lists (e.g., 80,443) and ranges (e.g., 1-1000)",
        "• Reports port status for each scanned port"
    ])
    
    # Slide 4: Features Part 2
    add_content_slide(prs, "Key Features (Part 2)", [
        "ARP Scan: Discovers active hosts on a local network segment",
        "• Returns IP and MAC address pairs",
        "• Useful for subnet mapping",
        "",
        "Flexible Targeting:",
        "• Scans single IPs: 192.168.1.1",
        "• Scans entire subnets: 192.168.1.0/24"
    ])
    
    # Slide 5: Features Part 3
    add_content_slide(prs, "Key Features (Part 3)", [
        "Performance & Control:",
        "• Timeout Control: User-configurable timeout per scan (default: 2 seconds)",
        "• Rate Limiting: Optional delay between probes to avoid flooding (--delay flag)",
        "• Max Ports Enforcement: Limits the number of ports scanned (--max-ports flag)"
    ])
    
    # Slide 6: Output & Reporting
    add_content_slide(prs, "Output & Reporting", [
        "• Terminal Display: Real-time results printed to console",
        "• File Output: Save results to text files (--output flag)",
        "• Metadata Headers: Includes scan timestamp, target, scan type, and parameters",
        "• Safe File Handling:",
        "  - Append mode (--append): Add to existing files",
        "  - Force overwrite (--force): Override safety checks",
        "  - Prevents accidental data loss"
    ])
    
    # Slide 7: Security Features
    add_content_slide(prs, "Security Features", [
        "✓ Privilege Checks: Warns if root privileges required",
        "✓ Rate Limiting: Avoids network flooding and DoS-like behavior",
        "✓ Logging: All scan activities logged to logs/scanner.log",
        "✓ Input Validation: Validates IP addresses, subnets, and port ranges",
        "✓ Safe Output Handling: Prevents accidental file overwrites"
    ])
    
    # Slide 8: Architecture
    add_content_slide(prs, "Project Architecture", [
        "CLI Interface (main.py)",
        "↓",
        "Argument Parsing & Validation",
        "↓",
        "Scanner Logic (scanner.py) → ICMP, TCP, ARP Scans",
        "↓",
        "Result Formatting → Terminal Print or File Write"
    ])
    
    # Slide 9: Technology Stack
    add_content_slide(prs, "Technology Stack", [
        "Language: Python 3.10+",
        "",
        "Core Libraries:",
        "• Scapy: Network packet manipulation and analysis",
        "• ipaddress: IP address validation and subnet handling",
        "• typing-extensions: Type hints for better code clarity",
        "• logging: Audit and error tracking"
    ])
    
    # Slide 10: Usage Examples
    add_code_slide(prs, "Usage Examples", 
                   "# Scan single host for ICMP and TCP ports\n"
                   "python src/main.py 192.168.1.1 -p 80,443\n\n"
                   "# Scan subnet with ARP and save results\n"
                   "python src/main.py 192.168.1.0/24 -t arp -o results.txt\n\n"
                   "# Full scan with rate limiting\n"
                   "python src/main.py 192.168.1.1 -t all -p 1-100 --delay 0.5 -o scan.txt")
    
    # Slide 11: CLI Options
    add_content_slide(prs, "Command-Line Options", [
        "-t, --type: Scan type (all, icmp, tcp, arp). Default: all",
        "-p, --ports: Ports to scan (e.g., 80,443 or 1-1000)",
        "-T, --timeout: Timeout in seconds. Default: 2",
        "-o, --output: Write output to text file",
        "--delay: Delay between probes in seconds. Default: 0.0",
        "--max-ports: Maximum ports per scan. Default: 1024",
        "--append: Append to output file instead of overwriting",
        "--force: Force actions that would otherwise be blocked"
    ])
    
    # Slide 12: Project Structure
    add_content_slide(prs, "Project Structure", [
        "Network Scanning /",
        "├── main.py (Entry point)",
        "├── src/scanner.py (Core scanning logic)",
        "├── tests/ (Unit and integration tests)",
        "├── docs/ (Documentation)",
        "├── requirements.txt (Dependencies)",
        "└── README.md (Quickstart guide)"
    ])
    
    # Slide 13: Sample Output
    add_code_slide(prs, "Sample Output",
                   "ICMP Scan Result:\n"
                   "Scan Results:\n"
                   "=============\n"
                   "ICMP Scan: Host 192.168.1.1 is UP\n\n"
                   "TCP Scan Result:\n"
                   "Port 80: open\n"
                   "Port 443: open\n"
                   "Port 8080: closed")
    
    # Slide 14: Advantages
    add_content_slide(prs, "Advantages & Benefits", [
        "✓ Comprehensive: Supports multiple scan types",
        "✓ User-Friendly: Simple CLI with clear help messages",
        "✓ Secure: Built-in privilege checks and logging",
        "✓ Efficient: Rate limiting to avoid network strain",
        "✓ Flexible: Works with IPs and subnets",
        "✓ Safe Output: Prevents accidental file overwrites",
        "✓ Extensible: Modular design for new features"
    ])
    
    # Slide 15: Testing
    add_content_slide(prs, "Testing & Quality", [
        "• Unit Tests: tests/test_scanner.py",
        "• Integration Tests: tests/test_integration.py",
        "• Comprehensive error handling and logging",
        "• All inputs validated before processing",
        "",
        "Run tests:",
        "python -m unittest discover tests"
    ])
    
    # Slide 16: Security Considerations
    add_content_slide(prs, "Security Considerations", [
        "⚠ Requires root/sudo privileges for raw socket operations",
        "⚠ Use responsibly and only on networks you own or have permission to scan",
        "⚠ Unauthorized network scanning may be illegal",
        "✓ Enable logging for audit trails and compliance",
        "✓ Rate limiting prevents accidental DoS-like behavior"
    ])
    
    # Slide 17: Future Enhancements
    add_content_slide(prs, "Future Enhancements", [
        "• UDP Scan: Add UDP port scanning capability",
        "• Service Detection: Identify services running on open ports",
        "• Web Dashboard: GUI for visualizing scan results",
        "• JSON Export: Output results in JSON format",
        "• Scheduling: Automated periodic scans",
        "• Database Integration: Store historical scan data"
    ])
    
    # Slide 18: Installation
    add_code_slide(prs, "Installation & Setup",
                   "# Create virtual environment\n"
                   "python3 -m venv venv\n\n"
                   "# Activate virtual environment\n"
                   "source venv/bin/activate\n\n"
                   "# Install dependencies\n"
                   "pip install -r requirements.txt\n\n"
                   "# Run a scan\n"
                   "python src/main.py 192.168.1.1 -p 80,443")
    
    # Slide 19: Documentation
    add_content_slide(prs, "Documentation", [
        "• README.md: Quick start guide and basic usage",
        "• Documentation_Overview.md: Detailed project documentation",
        "• Code Comments: Inline documentation in Python files",
        "• Help Command: python src/main.py --help for CLI reference"
    ])
    
    # Slide 20: Conclusion
    add_content_slide(prs, "Conclusion", [
        "What: Python network scanning tool with ICMP, TCP, and ARP",
        "Why: Essential for network administration and security auditing",
        "How: Simple CLI, flexible options, secure design",
        "",
        "Key Takeaways:",
        "✓ Multi-protocol scanning ✓ Flexible targeting",
        "✓ Secure implementation ✓ Safe output handling"
    ])
    
    # Save presentation
    output_path = "/home/ayush/Desktop/Network Scanning /Network_Scanning_Project.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
